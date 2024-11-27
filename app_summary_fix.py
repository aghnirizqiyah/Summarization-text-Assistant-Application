import streamlit as st
import nltk
from transformers import pipeline
from sumy.parsers.plaintext import PlaintextParser
from sumy.nlp.tokenizers import Tokenizer
from sumy.summarizers.luhn import LuhnSummarizer
import PyPDF2
import io
import docx
from pptx import Presentation

# Download required NLTK data
nltk.download('punkt')

def extract_text_from_pdf(file_bytes):
    """Extract text from PDF file"""
    pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def extract_text_from_docx(file_bytes):
    """Extract text from DOCX file"""
    doc = docx.Document(io.BytesIO(file_bytes))
    text = []
    for paragraph in doc.paragraphs:
        text.append(paragraph.text)
    return '\n'.join(text)

def extract_text_from_pptx(file_bytes):
    """Extract text from PPTX file"""
    presentation = Presentation(io.BytesIO(file_bytes))
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text.append(paragraph.text)
    return '\n'.join(text)

@st.cache_resource
def load_summarizer():
    """Load and cache the BART summarizer"""
    return pipeline("summarization", model="facebook/bart-large-cnn")

def abstractive_summary(text: str) -> str:
    """Generate abstractive summary using BART model"""
    summarizer = load_summarizer()
    max_chunk_length = 1024
    chunks = [text[i:i + max_chunk_length] for i in range(0, len(text), max_chunk_length)]
    
    summaries = []
    for chunk in chunks:
        summary = summarizer(chunk, max_length=130, min_length=30, do_sample=False)
        summaries.append(summary[0]['summary_text'])
    
    return ' '.join(summaries)

def sumy_summary(text: str, num_sentences: int = 5) -> str:
    """Generate extractive summary using Sumy (Luhn Algorithm)"""
    parser = PlaintextParser.from_string(text, Tokenizer("english"))
    summarizer = LuhnSummarizer()
    summary = summarizer(parser.document, num_sentences)
    return ' '.join(str(sentence) for sentence in summary)

def main():
    st.set_page_config(
        page_title="Summarization Assistant",
        page_icon="📄",
        layout="wide"
    )

    # Header Section
    st.markdown("""
        <style>
        body, .stMarkdown, .stButton button {
            color: black !important;
        }
        .css-1v3fvcr p, .css-1v3fvcr {
            color: black !important;
        }
        .main {
            background-color: #f4f4f4;
            padding: 20px;
            border-radius: 10px;
        }
        .custom-summary-text {
            font-size: 20px;
            line-height: 1.8;
            font-family: 'Arial', sans-serif;
            color: #333333;
        }
        </style>
        """, unsafe_allow_html=True)
    
    st.title("📄 Summarization Assistant")
    st.markdown("### **Intelligent Document Summarization with AI**")

    # Settings Section
    st.markdown("---")
    st.subheader("⚙️ Settings")
    col1, col2 = st.columns([1, 3])
    with col1:
        method = st.radio(
            "Summarization Method:",
            options=['Extraction', 'Abstractive', 'Both'],
            index=0,
            horizontal=True
        )
    with col2:
        file_type = st.selectbox(
            "Select Document Type:",
            options=['PDF Document (.pdf)', 'Text File (.txt)', 'Word Document (.docx)', 'PowerPoint Presentation (.pptx)']
        )

    # About Section
    st.markdown("---")
    st.markdown("### ℹ️ About")
    st.markdown("""
    **Summarization Methods:**

    - **Extraction Summarization**: This method picks the most important sentences directly from the original text using the SUMY (Luhn) algorithm, maintaining their original wording.
    - **Abstractive Summarization**: This method creates a completely new summary by rephrasing information in a more natural and concise language using the BART algorithm.
    - **Both Methods**: Displays summaries generated from both Abstractive and Extraction methods side-by-side for comparison.

    Choose the method that suits your needs:
    - Use **Extraction** for a direct and structured extraction.
    - Use **Abstractive** for a creative, rewritten summary.
    - Use **Both** to compare the strengths of both approaches!
""")


    st.markdown("---")

    # File Upload Section
    st.subheader("📤 Upload Your Document")
    uploaded_file = st.file_uploader(
        "Drag and drop a file here:",
        type={
            'Text File (.txt)': 'txt',
            'PDF Document (.pdf)': 'pdf',
            'Word Document (.docx)': 'docx',
            'PowerPoint Presentation (.pptx)': 'pptx'
        }[file_type],
        help=f"Upload a {file_type} document"
    )

    # Process Uploaded File
    if uploaded_file is not None:
        col1, col2 = st.columns([1, 2])
        with col1:
            st.write(f"**Filename:** {uploaded_file.name}")
            st.write(f"**File Size:** {uploaded_file.size / 1024:.2f} KB")

        if st.button("🔍 Generate Summary"):
            with st.spinner("🔄 Processing..."):
                try:
                    # Extract Text
                    if file_type == 'PDF Document (.pdf)':
                        text = extract_text_from_pdf(uploaded_file.read())
                    elif file_type == 'Word Document (.docx)':
                        text = extract_text_from_docx(uploaded_file.read())
                    elif file_type == 'PowerPoint Presentation (.pptx)':
                        text = extract_text_from_pptx(uploaded_file.read())
                    else:
                        text = uploaded_file.read().decode()

                    summaries = {}
                    if method in ["Extraction", "Both"]:
                        summaries["Extraction Summary"] = sumy_summary(text)
                    if method in ["Abstractive", "Both"]:
                        summaries["Abstractive Summary"] = abstractive_summary(text)

                    # Results Section
                    st.markdown("---")
                    st.header("📋 Results")
                    st.subheader("📊 Text Statistics")
                    stat_col1, stat_col2, stat_col3 = st.columns(3)
                    stat_col1.metric("Original Text Length", f"{len(text)} chars")

                    if method == "Both":
                        total_summary_length = sum(len(summary) for summary in summaries.values())
                        stat_col2.metric("Total Summary Length", f"{total_summary_length} chars")
                        stat_col3.metric("Avg Compression Ratio", f"{(total_summary_length / len(text)) * 100:.1f}%")
                    else:
                        selected_summary = list(summaries.values())[0]
                        stat_col2.metric("Summary Length", f"{len(selected_summary)} chars")
                        stat_col3.metric("Compression Ratio", f"{(len(selected_summary) / len(text)) * 100:.1f}%")

                    # Display Summaries in Grid Layout
                    st.markdown("---")
                    st.subheader("📜 Summaries")
                    tab1, tab2, tab3 = st.tabs(["Extraction Summary", "Abstractive Summary", "Original Text"])

                    if "Extraction Summary" in summaries:
                        with tab1:
                            st.markdown(f"""
                                <div class="custom-summary-text">
                                    {summaries["Extraction Summary"]}
                                </div>
                            """, unsafe_allow_html=True)
                            st.download_button(
                                label="Download Extraction Summary",
                                data=summaries["Extraction Summary"],
                                file_name=f"sumy_summary_{uploaded_file.name}.txt",
                                mime="text/plain"
                            )
                    if "Abstractive Summary" in summaries:
                        with tab2:
                            st.markdown(f"""
                                <div class="custom-summary-text">
                                    {summaries["Abstractive Summary"]}
                                </div>
                            """, unsafe_allow_html=True)
                            st.download_button(
                                label="Download Abstractive Summary",
                                data=summaries["Abstractive Summary"],
                                file_name=f"abstractive_summary_{uploaded_file.name}.txt",
                                mime="text/plain"
                            )

                    # Original Text in an Expander
                    with tab3:
                        with st.expander("📰 View Original Text"):
                            st.write(text)

                except Exception as e:
                    st.error(f"An error occurred: {str(e)}")



if __name__ == "__main__":
    main()
